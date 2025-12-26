import os
import sqlite3
import base64
import requests
import matplotlib.pyplot as plt

from fastapi import FastAPI, Request, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, FileResponse, JSONResponse

from slowapi import Limiter
from slowapi.util import get_remote_address

from groq import Groq
from fpdf import FPDF
from pptx import Presentation
import replicate

# ==========================================================
# CONFIG
# ==========================================================
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
NEWS_API_KEY = os.getenv("NEWS_API_KEY")
REPLICATE_API_TOKEN = os.getenv("REPLICATE_API_TOKEN")

client = Groq(api_key=GROQ_API_KEY)
replicate.Client(api_token=REPLICATE_API_TOKEN)

app = FastAPI()
limiter = Limiter(key_func=get_remote_address)
app.state.limiter = limiter

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ==========================================================
# SYSTEM PROMPT (INTELLIGENCE BOOST)
# ==========================================================
SYSTEM_PROMPT = """
You are lEvO — an advanced AI assistant.

Capabilities:
- Think deeply and reason step-by-step
- Give long, structured, detailed answers
- Use headings, bullet points, tables
- Be accurate and honest
- Always output clickable links

Style:
- Friendly, confident, modern
- Not robotic
- Helpful and intelligent

If live data is provided, prioritize it.
"""

# ==========================================================
# MEMORY (SQLite)
# ==========================================================
conn = sqlite3.connect("memory.db", check_same_thread=False)
cursor = conn.cursor()
cursor.execute("CREATE TABLE IF NOT EXISTS memory (user_id TEXT, text TEXT)")
conn.commit()

def get_memory(uid):
    cursor.execute("SELECT text FROM memory WHERE user_id=?", (uid,))
    return "\n".join([row[0] for row in cursor.fetchall()])

def save_memory(uid, msg):
    cursor.execute("INSERT INTO memory VALUES (?, ?)", (uid, msg))
    conn.commit()

# ==========================================================
# HOME
# ==========================================================
@app.get("/")
def home():
    return {"status": "lEvO AI API is LIVE 🔥"}

# ==========================================================
# LIVE NEWS
# ==========================================================
@app.get("/news")
def news(topic: str = "technology"):
    url = f"https://newsapi.org/v2/top-headlines?q={topic}&apiKey={NEWS_API_KEY}"
    return requests.get(url).json()

# ==========================================================
# NORMAL CHAT (LONG RESPONSES)
# ==========================================================
@app.post("/chat")
@limiter.limit("15/minute")
async def chat(request: Request):
    data = await request.json()
    prompt = data.get("prompt", "")
    user_id = data.get("user_id", "guest")

    memory = get_memory(user_id)

    msgs = [
        {"role": "system", "content": SYSTEM_PROMPT},
        {"role": "system", "content": f"Conversation memory:\n{memory}"},
        {"role": "user", "content": prompt}
    ]

    completion = client.chat.completions.create(
        model="llama-3.1-8b-instant",
        messages=msgs,
        max_tokens=4096
    )

    reply = completion.choices[0].message.content
    save_memory(user_id, prompt)
    save_memory(user_id, reply)

    return {"response": reply}

# ==========================================================
# STREAMING CHAT
# ==========================================================
@app.get("/stream")
async def stream(prompt: str):
    def generate():
        completion = client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[
                {"role": "system", "content": SYSTEM_PROMPT},
                {"role": "user", "content": prompt}
            ],
            stream=True
        )
        for chunk in completion:
            if chunk.choices and chunk.choices[0].delta.content:
                yield f"data: {chunk.choices[0].delta.content}\n\n"
        yield "data: [DONE]\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")

# ==========================================================
# IMAGE ANALYSIS
# ==========================================================
@app.post("/vision")
async def vision(file: UploadFile = File(...)):
    img_bytes = await file.read()
    img_b64 = base64.b64encode(img_bytes).decode()

    completion = client.chat.completions.create(
        model="llava-v1.6",
        messages=[
            {"role": "user", "content": [
                {"type": "input_text", "text": "Analyze this image deeply."},
                {"type": "input_image", "image": img_b64},
            ]}
        ]
    )

    return {"response": completion.choices[0].message.content}

# ==========================================================
# IMAGE GENERATION (SDXL)
# ==========================================================
@app.post("/image-generate")
def image_generate(prompt: str):
    output = replicate.run(
        "stability-ai/sdxl",
        input={"prompt": prompt}
    )
    return {"image_url": output[0]}

# ==========================================================
# PDF CREATION
# ==========================================================
@app.post("/create-pdf")
def create_pdf(text: str):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(0, 8, text)
    pdf.output("output.pdf")
    return FileResponse("output.pdf")

# ==========================================================
# TXT FILE
# ==========================================================
@app.post("/create-txt")
def create_txt(text: str):
    with open("output.txt", "w") as f:
        f.write(text)
    return FileResponse("output.txt")

# ==========================================================
# PPTX
# ==========================================================
@app.post("/create-ppt")
def create_ppt(text: str):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Generated by lEvO"
    slide.placeholders[1].text = text
    prs.save("output.pptx")
    return FileResponse("output.pptx")

# ==========================================================
# GRAPH
# ==========================================================
@app.get("/graph")
def graph():
    plt.plot([1,2,3,4], [10,20,15,30])
    plt.savefig("graph.png")
    return FileResponse("graph.png")

# ==========================================================
# GOOGLE SEARCH (RAG)
# ==========================================================
@app.get("/search")
def search(q: str):
    url = f"https://ddg-api.herokuapp.com/search?query={q}"
    return requests.get(url).json()
